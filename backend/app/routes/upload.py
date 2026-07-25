import asyncio
import json
import logging
import os
import re
import uuid
from pathlib import Path
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.orm import Session
from app.dependencies import get_db
from app.models import Article
from app.schemas import ArticleResponse
from app.config import (
    LLM_API_KEY, LLM_API_BASE, LLM_MODEL,
    VISION_API_KEY, VISION_API_BASE, VISION_MODEL,
    ASR_API_KEY, ASR_API_BASE, ASR_MODEL,
    UPLOAD_DIR as UPLOAD_DIR_STR,
)
from app.utils import find_ffmpeg

router = APIRouter(prefix="/api/upload", tags=["upload"])

logger = logging.getLogger(__name__)

# ─── Config ────────────────────────────────────────

UPLOAD_DIR = Path(UPLOAD_DIR_STR)

# Ensure upload directory exists
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

# Supported formats
TEXT_EXTENSIONS = {'.txt', '.md', '.markdown', '.json', '.xml', '.csv', '.yaml', '.yml', '.py', '.js', '.ts', '.html', '.css'}
WORD_EXTENSIONS = {'.docx'}
EXCEL_EXTENSIONS = {'.xlsx', '.xls'}
PPT_EXTENSIONS = {'.pptx', '.ppt'}
PDF_EXTENSIONS = {'.pdf'}
AUDIO_EXTENSIONS = {'.mp3', '.wav', '.m4a', '.flac', '.ogg', '.wma'}
VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov', '.mkv', '.webm', '.wmv'}
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.svg', '.bmp', '.ico', '.tiff', '.tif'}


# ─── File parsers ──────────────────────────────────

async def parse_text(file: UploadFile) -> str:
    """Parse plain text files."""
    content = await file.read()
    return content.decode("utf-8", errors="replace")


def parse_text_from_bytes(content_bytes: bytes) -> str:
    """Parse plain text from already-read bytes."""
    return content_bytes.decode("utf-8", errors="replace")


def parse_docx(file_path: str) -> str:
    """Parse Word documents."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def parse_xlsx(file_path: str) -> str:
    """Parse Excel files — iterate all sheets, join cell values."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            text = " | ".join(cells).strip()
            if text:
                rows.append(text)
        parts.append("\n".join(rows))
    return "\n\n".join(parts)


def parse_pptx(file_path: str) -> str:
    """Parse PowerPoint files."""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"## Slide {i + 1}\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def parse_pdf(file_path: str) -> str:
    """Parse PDF files."""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts)


async def parse_image(file_path: str, original_name: str) -> str:
    """Describe image content via vision LLM.
    Returns markdown with an embedded image + description."""
    import base64
    import httpx

    storage_name = Path(file_path).name
    img_tag = f'<img src="/api/media/{storage_name}" alt="{original_name}" style="max-width:100%;height:auto;display:block;border-radius:4px">'

    if not VISION_API_KEY:
        return f"{img_tag}\n\n# 图片：{original_name}\n\n> 未配置视觉模型 API，无法自动描述图片内容。请手动添加。"

    # Read and encode image
    with open(file_path, "rb") as f:
        image_data = base64.b64encode(f.read()).decode("utf-8")

    # Determine MIME type
    ext = Path(original_name).suffix.lower()
    mime_map = {
        '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
        '.gif': 'image/gif', '.webp': 'image/webp', '.svg': 'image/svg+xml',
        '.bmp': 'image/bmp', '.ico': 'image/x-icon', '.tiff': 'image/tiff', '.tif': 'image/tiff',
    }
    mime = mime_map.get(ext, 'image/png')

    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(
                f"{VISION_API_BASE.rstrip('/')}/chat/completions",
                headers={
                    "Authorization": f"Bearer {VISION_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": VISION_MODEL,   # vision model for image understanding
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "请用中文详细描述这张图片的内容（不超过300字）。如果是文档截图、表格、图表，请尽可能提取其中的文字信息。",
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:{mime};base64,{image_data}",
                                    },
                                },
                            ],
                        }
                    ],
                    "max_tokens": 800,
                },
            )
            if resp.status_code == 200:
                data = resp.json()
                desc = data["choices"][0]["message"]["content"].strip()
                return f"{img_tag}\n\n# 图片描述：{original_name}\n\n{desc}"
            else:
                logger.warning("Image description API returned %d: %s", resp.status_code, resp.text[:200])
    except Exception as e:
        logger.warning("Image description failed: %s", e)

    return f"{img_tag}\n\n# 图片：{original_name}\n\n> 图片描述生成失败。请手动添加文章内容。"


async def parse_media(content_bytes: bytes, filename: str, content_type: str) -> str:
    """Handle audio — try ASR transcription, surface errors in content."""
    name_no_ext = Path(filename).stem.replace('_', ' ').replace('-', ' ')

    logger.info(f"[AUDIO] parse_media called: filename={filename}, size={len(content_bytes)}, ASR_API_KEY={'set' if ASR_API_KEY else 'NOT SET'}")
    if ASR_API_KEY:
        try:
            logger.info("[AUDIO] Calling transcribe_media_from_bytes...")
            result = await transcribe_media_from_bytes(content_bytes, filename, content_type)
            if result and result.strip():
                # Check if it's an error message
                if result.startswith("ERROR:"):
                    err_msg = result[len("ERROR:"):].strip()
                    return (
                        f"# 音频：{name_no_ext}\n\n"
                        f"> ⚠️ 语音识别失败：{err_msg}\n\n"
                        f"> 文件名：{filename}\n"
                        f"> 类型：{content_type}"
                    )
                return result
        except Exception:
            pass

    # Fallback
    return (
        f"# 音频：{name_no_ext}\n\n"
        f"该音频文件记录了{name_no_ext}相关的内容。\n\n"
        f"> 文件名：{filename}\n"
        f"> 类型：{content_type}\n"
        f"> 注意：无法自动转写此音频的内容。请手动添加描述。"
    )


async def parse_video(file_path: str, original_name: str) -> str:
    """Extract key frames from video and describe via vision LLM.
    Returns markdown with an embedded video player + description."""
    import base64
    import cv2
    import httpx

    storage_name = Path(file_path).name
    video_tag = f'<video controls src="/api/media/{storage_name}" style="width:100%;max-width:100%"></video>'
    name_no_ext = Path(original_name).stem.replace('_', ' ').replace('-', ' ')

    if not VISION_API_KEY:
        return f"{video_tag}\n\n# 视频：{name_no_ext}\n\n> 未配置视觉模型 API，无法自动描述视频内容。"

    # ── Extract key frames ──
    cap = cv2.VideoCapture(file_path)
    if not cap.isOpened():
        return f"{video_tag}\n\n# 视频：{name_no_ext}\n\n> 无法打开视频文件。"

    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    fps = cap.get(cv2.CAP_PROP_FPS)
    duration = total_frames / fps if fps > 0 else 0

    # Extract up to 5 frames at 0%, 25%, 50%, 75%, 90% of the video
    positions = [0, 0.25, 0.5, 0.75, 0.9]
    frames_b64: list[str] = []

    for pos in positions:
        frame_idx = int(total_frames * pos)
        if frame_idx >= total_frames:
            frame_idx = total_frames - 1
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
        ret, frame = cap.read()
        if ret and frame is not None:
            # Resize large frames to max 1024px on longest side (API size limits)
            h, w = frame.shape[:2]
            max_side = max(h, w)
            if max_side > 1024:
                scale = 1024 / max_side
                frame = cv2.resize(frame, (int(w * scale), int(h * scale)))
            # Encode as JPEG
            _, buf = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 80])
            frames_b64.append(base64.b64encode(buf).decode("utf-8"))

    cap.release()

    if not frames_b64:
        return f"{video_tag}\n\n# 视频：{name_no_ext}\n\n> 无法从视频中提取画面。"

    # ── Send frames to vision model ──
    n_frames = len(frames_b64)
    user_content: list[dict] = [
        {
            "type": "text",
            "text": (
                f"这是一个视频文件「{name_no_ext}」，时长约 {duration:.0f} 秒。"
                f"以下是从视频中提取的 {n_frames} 个关键帧画面（按时间顺序）。"
                "请用中文详细描述这个视频的内容（不超过 300 字），包括：\n"
                "1. 视频的主题和场景\n"
                "2. 出现的人物、物体或活动\n"
                "3. 画面随时间的变化\n"
                "4. 视频想要传达的信息"
            ),
        }
    ]
    for i, b64 in enumerate(frames_b64):
        user_content.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/jpeg;base64,{b64}"},
        })

    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(
                f"{VISION_API_BASE.rstrip('/')}/chat/completions",
                headers={
                    "Authorization": f"Bearer {VISION_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": VISION_MODEL,
                    "messages": [
                        {
                            "role": "user",
                            "content": user_content,
                        }
                    ],
                    "max_tokens": 800,
                },
            )
            if resp.status_code == 200:
                data = resp.json()
                desc = data["choices"][0]["message"]["content"].strip()
                logger.info(f"[VIDEO] Vision model response ({len(desc)} chars): {desc[:150]}")
                return f"{video_tag}\n\n# 视频内容描述：{name_no_ext}\n\n{desc}"
            else:
                logger.info(f"[VIDEO] Vision API returned {resp.status_code}: {resp.text[:200]}")
    except Exception as e:
        logger.info(f"[VIDEO] Vision API error: {e}")

    # Fallback: if vision model fails
    return (
        f"{video_tag}\n\n"
        f"# 视频：{name_no_ext}\n\n"
        f"该视频文件记录了{name_no_ext}相关的内容（时长约 {duration:.0f} 秒）。\n\n"
        f"> 视频内容自动识别失败。请手动添加描述。"
    )


def _convert_to_mono_wav(audio_bytes: bytes, orig_ext: str) -> bytes | None:
    """Convert audio to 16kHz mono WAV bytes. Returns None if conversion fails."""
    import io
    import subprocess
    import wave
    import audioop

    ext = orig_ext.lower()

    # ── WAV: use built-in wave module (no ffmpeg needed) ──
    if ext == '.wav':
        try:
            with wave.open(io.BytesIO(audio_bytes), 'rb') as wf:
                nchannels = wf.getnchannels()
                sampwidth = wf.getsampwidth()
                framerate = wf.getframerate()
                frames = wf.readframes(wf.getnframes())

            # Convert to mono if needed
            if nchannels > 1:
                frames = audioop.tomono(frames, sampwidth, 1.0, 1.0)
                logger.info(f"[AUDIO] Converted {nchannels}ch → mono, {sampwidth*8}bit")

            # Resample to 16kHz if needed
            if framerate != 16000:
                frames = audioop.ratecv(frames, sampwidth, 1, framerate, 16000, None)[0]
                logger.info(f"[AUDIO] Resampled {framerate}Hz → 16000Hz")

            # Write back as mono WAV
            buf = io.BytesIO()
            with wave.open(buf, 'wb') as wf:
                wf.setnchannels(1)
                wf.setsampwidth(sampwidth)
                wf.setframerate(16000)
                wf.writeframes(frames)
            return buf.getvalue()
        except Exception as e:
            logger.info(f"[AUDIO] WAV conversion failed: {e}")
            return None

    # ── Other formats: use ffmpeg subprocess ──
    ffmpeg_path = find_ffmpeg()
    if not ffmpeg_path:
        logger.info("[AUDIO] ffmpeg not found — cannot convert non-WAV audio")
        return None

    try:
        logger.info(f"[AUDIO] Converting {ext} to mono 16kHz WAV via ffmpeg: {ffmpeg_path}")
        result = subprocess.run(
            [
                ffmpeg_path,
                '-i', 'pipe:0',       # read from stdin
                '-ac', '1',            # mono
                '-ar', '16000',        # 16kHz
                '-f', 'wav',           # WAV output
                'pipe:1',              # write to stdout
            ],
            input=audio_bytes,
            capture_output=True,
            timeout=60,
        )
        if result.returncode != 0:
            stderr = result.stderr.decode('utf-8', errors='replace')[:300]
            logger.info(f"[AUDIO] ffmpeg error: {stderr}")
            return None

        logger.info(f"[AUDIO] ffmpeg conversion OK, output size={len(result.stdout)}")
        return result.stdout
    except FileNotFoundError:
        logger.info("[AUDIO] ffmpeg executable not found at path")
        return None
    except Exception as e:
        logger.info(f"[AUDIO] ffmpeg conversion failed: {e}")
        return None


async def transcribe_media_from_bytes(content: bytes, filename: str, content_type: str) -> str:
    """Transcribe audio via Zhipu ASR model (GLM-ASR-2512).

    Returns the transcribed text on success, or an error message string
    prefixed with "ERROR:" on failure (so the caller can surface it).
    """
    import httpx

    ext = Path(filename).suffix.lower()
    logger.info(f"[AUDIO] transcribe_media_from_bytes: filename={filename}, ext={ext}, size={len(content)}")

    # ── Convert to mono 16kHz WAV (required by GLM-ASR-2512) ──
    logger.info("[AUDIO] Calling _convert_to_mono_wav...")
    mono_bytes = _convert_to_mono_wav(content, ext)
    if mono_bytes is None:
        return "ERROR: 音频格式转换失败（需要单声道音频）。请尝试转换音频文件后重新上传。"

    # ── Zhipu ASR endpoint ──
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(
                f"{ASR_API_BASE.rstrip('/').replace('/chat/completions', '')}/audio/transcriptions",
                headers={"Authorization": f"Bearer {ASR_API_KEY}"},
                files={"file": ("audio.wav", mono_bytes, "audio/wav")},
                data={"model": ASR_MODEL},
            )
            if resp.status_code == 200:
                data = resp.json()
                text = data.get("text", "")
                if text and text.strip():
                    logger.info(f"[ASR] GLM-ASR-2512 success: {text[:150]}")
                    return f"# 音频转录\n\n{text}"
                return "ERROR: GLM-ASR-2512 返回了空文本。"
            else:
                err_detail = resp.text[:300]
                logger.info(f"[ASR] GLM-ASR-2512 returned {resp.status_code}: {err_detail}")
                return f"ERROR: GLM-ASR-2512 识别失败（HTTP {resp.status_code}）：{err_detail}"
    except Exception as e:
        logger.info(f"[ASR] GLM-ASR-2512 error: {e}")
        return f"ERROR: GLM-ASR-2512 调用异常：{e}"


# ─── LLM helpers ───────────────────────────────────

def _clean_text_for_title(text: str) -> str:
    """Strip HTML tags, markdown images, and metadata headers — keep the real content."""
    import re as _re
    cleaned = _re.sub(r'<[^>]+>', '', text)                # HTML tags
    cleaned = _re.sub(r'!\[[^\]]*\]\([^)]+\)', '', cleaned) # markdown images
    # Remove metadata headers produced by our parsers
    cleaned = _re.sub(r'^#+\s*图片描述[：:].*\n?', '', cleaned, flags=_re.MULTILINE)
    cleaned = _re.sub(r'^#+\s*图片[：:].*\n?', '', cleaned, flags=_re.MULTILINE)
    cleaned = _re.sub(r'^#+\s*文件内容描述\s*\n?', '', cleaned, flags=_re.MULTILINE)
    cleaned = _re.sub(r'^#+\s*音视频文件\s*\n?', '', cleaned, flags=_re.MULTILINE)
    # Remove metadata lines
    cleaned = _re.sub(r'^-\s*(?:文件名|类型|大小)[：:].*\n?', '', cleaned, flags=_re.MULTILINE)
    cleaned = _re.sub(r'^>.*\n?', '', cleaned, flags=_re.MULTILINE)
    cleaned = _re.sub(r'\n{3,}', '\n\n', cleaned)
    return cleaned.strip()


async def generate_title(text: str) -> str:
    """Generate a concise summary title from document content via LLM.

    Returns a title that *synthesizes* the document's core topic (not a sentence
    copied from the text).  Returns empty string on failure so the caller can
    decide the fallback strategy.
    """
    if not LLM_API_KEY:
        logger.info("[TITLE] LLM not configured — skipping title generation")
        return ""

    import httpx

    # Use cleaned text — strip media tags and metadata headers
    cleaned = _clean_text_for_title(text)
    logger.info(f"[TITLE] Cleaned text: {len(cleaned)} chars, first 100: {cleaned[:100]}")

    if len(cleaned) < 10:
        logger.info(f"[TITLE] Text too short ({len(cleaned)} chars) — skipping")
        return ""

    # Send more context for better understanding (up to 8000 chars)
    context = cleaned[:8000]

    prompt = (
        "你是一个专业的文档摘要专家。请仔细阅读以下文档内容，在**理解全文主旨**的基础上，"
        "用一句完整、精炼的话概括文档的核心内容，作为标题。\n\n"
        "核心原则：\n"
        "1. 先通读全文，理解文档在讲什么，然后用自己的话提炼标题\n"
        "2. 标题应该概括文档的整体主题，而不是照抄文中的某一句\n"
        "3. 让读者一眼就能知道这份文档是关于什么的\n\n"
        "具体要求：\n"
        "- 15-40 个字，简洁完整\n"
        "- 技术文档 → 概括技术主题和要点\n"
        "- 聊天记录/对话 → 概括讨论的主要话题和结论\n"
        "- 图片描述 → 概括图片的主要内容和场景\n"
        "- 音视频内容 → 概括主题和关键信息\n"
        "- 不要包含\"标题：\"\"本文\"\"该文档\"等冗余词语\n"
        "- 不要使用引号、书名号或 Markdown 格式\n"
        "- 只输出标题文本本身，不要任何解释\n\n"
        "好的标题示例：\n"
        "- Python 异步编程的核心概念与最佳实践\n"
        "- 2025 年产品路线图与关键里程碑规划\n"
        "- 基于深度学习的图像分类模型优化方法\n"
        "- 团队关于微服务架构迁移的技术方案讨论\n"
        "- 城市夜景航拍照片，展示 CBD 核心区的灯光与建筑群\n\n"
        "差的标题（避免）：\n"
        "- 在本文中我们将讨论异步编程  ← 冗余前缀\n"
        "- 第一章 概述  ← 没有实际信息\n"
        "- 接下来我们来看一下  ← 口语化，无概括\n\n"
        f"---\n文档内容：\n\n{context}\n\n---\n"
        "请为以上文档生成标题（只输出标题本身）："
    )

    try:
        logger.info(f"[TITLE] Calling LLM model={LLM_MODEL} with {len(context)} chars of context...")
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(
                f"{LLM_API_BASE.rstrip('/')}/chat/completions",
                headers={
                    "Authorization": f"Bearer {LLM_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": LLM_MODEL,
                    "messages": [
                        {"role": "user", "content": prompt},
                    ],
                    "max_tokens": 4000,  # GLM-5.2 reasoning model needs headroom for thinking
                    "temperature": 0.7,
                },
            )
            resp.raise_for_status()
            data = resp.json()
            logger.info("[TITLE] API response keys: {list(data.keys())}, choices={len(data.get('choices', []))}")
            if data.get("choices"):
                c0 = data["choices"][0]
                logger.info("[TITLE] Choice[0]: finish_reason={c0.get('finish_reason')}, message={c0.get('message')}")
            raw = (data.get("choices", [{}])[0].get("message", {}).get("content", "") or "").strip()
            logger.info(f"[TITLE] LLM raw response ({len(raw)} chars): {raw[:200]}")

            # Clean up formatting
            title = raw
            title = re.sub(r'^["\'"\'「『【《〈』」』】》〉]', '', title)
            title = re.sub(r'["\'"\'「『【《〈』」』】》〉]$', '', title)
            title = re.sub(r'^(?:标题|题目)[：:]\s*', '', title)
            title = re.sub(r'^#+\s*', '', title)
            title = re.sub(r'\n.*', '', title)             # first line only
            title = title.strip()

            # Validate: must be a meaningful phrase (not just a heading number or filler)
            if not title or len(title) < 4:
                logger.info(f"[TITLE] Title too short ({len(title) if title else 0} chars), discarding")
                return ""

            # Reject responses that are obviously not titles
            no_title_patterns = [
                r'^[第序]\s*\d+\s*[章节篇]',       # "第一章", "第3节"
                r'^[\(（]\s*[\)）]\s*$',            # just "()"
                r'^[一二三四五六七八九十]、',         # "一、概述"
                r'^\(?\d+\)[\.、]',                  # "1.", "1、"
                r'^(?:好的|以下|这里|下面是|例如)',    # meta-language
            ]
            for pat in no_title_patterns:
                if re.match(pat, title):
                    logger.info(f"[TITLE] Rejected meta/noise title: {title}")
                    return ""

            if len(title) > 100:
                title = title[:100]

            logger.info(f"[TITLE] ✅ Final title: {title}")
            return title

    except Exception as e:
        logger.info(f"[TITLE] ❌ Exception: {e}")
        return ""



async def extract_entities_and_relations(text: str) -> tuple[list[str], dict | None]:
    """Extract tags and structured entities+relations from document text via LLM.
    Returns (tags, entities_dict).
    Uses the shared extraction function via asyncio.to_thread to avoid blocking."""
    import asyncio
    from app.llm_extract import extract_tags_and_entities
    return await asyncio.to_thread(extract_tags_and_entities, text)


# ─── Route ─────────────────────────────────────────

@router.post("", response_model=ArticleResponse, status_code=201)
async def upload_file(
    file: UploadFile = File(...),
    category_id: str | None = Form(default=None),
    db: Session = Depends(get_db),
):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    # Determine file extension
    ext = Path(file.filename).suffix.lower()

    # 1. Save file to disk (sanitize filename to prevent path traversal)
    file_id = str(uuid.uuid4())
    safe_filename = re.sub(r'[^\w.\-]', '_', Path(file.filename).name)
    safe_name = f"{file_id}_{safe_filename}"
    file_path = UPLOAD_DIR / safe_name

    content_bytes = await file.read()

    # Prevent resource exhaustion: limit to 500MB
    if len(content_bytes) > 500 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (max 500MB)")

    with open(file_path, "wb") as f:
        f.write(content_bytes)

    # 2. Build initial content (no LLM — immediate display of media)
    is_media = ext in IMAGE_EXTENSIONS or ext in AUDIO_EXTENSIONS or ext in VIDEO_EXTENSIONS
    is_doc = ext in (TEXT_EXTENSIONS | WORD_EXTENSIONS | EXCEL_EXTENSIONS | PPT_EXTENSIONS | PDF_EXTENSIONS)
    media_src = f"/api/media/{safe_name}"
    raw_text = ""  # text extracted without LLM help

    try:
        if ext in TEXT_EXTENSIONS:
            raw_text = parse_text_from_bytes(content_bytes)
        elif ext in WORD_EXTENSIONS:
            raw_text = await asyncio.to_thread(parse_docx, str(file_path))
        elif ext in EXCEL_EXTENSIONS and ext != '.csv':
            raw_text = await asyncio.to_thread(parse_xlsx, str(file_path))
        elif ext in PPT_EXTENSIONS:
            raw_text = await asyncio.to_thread(parse_pptx, str(file_path))
        elif ext in PDF_EXTENSIONS:
            raw_text = await asyncio.to_thread(parse_pdf, str(file_path))
        elif ext in IMAGE_EXTENSIONS:
            raw_text = f'<img src="{media_src}" alt="{file.filename}" style="max-width:100%;height:auto;display:block;border-radius:4px">'
        elif ext in AUDIO_EXTENSIONS:
            raw_text = f'<audio controls src="{media_src}" style="width:100%"></audio>'
        elif ext in VIDEO_EXTENSIONS:
            raw_text = f'<video controls src="{media_src}" style="width:100%"></video>'
        else:
            raw_text = f"# {file.filename}\n\n不支持的文件格式。文件已作为附件保存。"
    except Exception as e:
        raw_text = f"# {file.filename}\n\n文件解析失败，文件已作为附件保存。"

    # 3. Add persistent doc-attachment marker for non-media files (so the frontend
    #    can uniquely identify each file, even when multiple share the same name).
    if is_doc:
        raw_text += f"\n\n<!-- doc-attachment: {file.filename} | {safe_name} -->"

    # 4. Track upload order so later-added attachments preserve correct ordering
    raw_text += f"\n\n<!-- attachments-order: {file.filename} -->"

    # 5. Use filename (without extension) as initial title; background task will generate a better one
    title = Path(file.filename).stem or file.filename

    # 5. Create article immediately — media visible, marked as processing
    article = Article(
        title=title,
        content=raw_text,
        category_id=category_id or None,
        tags=json.dumps([], ensure_ascii=False),
        entities=None,
        processing="processing",
        attachment_path=str(safe_name),
        attachment_name=file.filename,
        attachment_type=file.content_type or "",
    )
    db.add(article)
    db.commit()
    db.refresh(article)

    # 5. Background: LLM enhance title + content + tags + entities
    article_id = article.id

    async def _bg_enhance():
        from app.database import SessionLocal
        db2 = SessionLocal()
        try:
            # Step A: Generate full content with LLM description (for media)
            if is_media:
                if ext in IMAGE_EXTENSIONS:
                    full_text = await parse_image(str(file_path), file.filename)
                elif ext in VIDEO_EXTENSIONS:
                    full_text = await parse_video(str(file_path), file.filename)
                else:
                    # Re-read from disk to avoid capturing content_bytes in closure
                    with open(file_path, "rb") as f:
                        audio_bytes = f.read()
                    desc = await parse_media(audio_bytes, file.filename, file.content_type or "")
                    full_text = raw_text + "\n\n" + desc
            else:
                full_text = raw_text

            # Step B: Extract tags/entities (title is already set from filename)
            bg_tags, bg_entities = await extract_entities_and_relations(full_text)

            art = db2.query(Article).filter(Article.id == article_id).first()
            if not art:
                return

            # Collect errors for transparency
            errs: list[str] = []

            # Update content with full description
            if is_media:
                art.content = full_text

            # Update tags + entities; report if LLM extraction returned nothing
            if bg_tags:
                art.tags = json.dumps(bg_tags, ensure_ascii=False)
            else:
                errs.append("标签提取未返回结果")

            if bg_entities:
                art.entities = json.dumps(bg_entities, ensure_ascii=False)
            else:
                errs.append("实体和关系提取未返回结果")

            # Append error notes to content so the user can see what happened
            if errs:
                err_lines = "\n".join(f"- {e}" for e in errs)
                art.content = (art.content or raw_text) + f"\n\n> ⚠️ 以下步骤未成功完成：\n> \n> {err_lines}\n>\n> 模型: {LLM_MODEL} / {VISION_MODEL}"

            art.processing = None  # mark as done

            # Compute embeddings now that content has been enriched by LLM recognition
            try:
                from app.routes.qa import chunk_article, get_embedding
                from app.models import ArticleChunk

                chunks = chunk_article(art.content or "")
                for i, chunk_text in enumerate(chunks):
                    try:
                        vec = await get_embedding(chunk_text)
                        db2.add(ArticleChunk(
                            article_id=art.id,
                            chunk_index=str(i),
                            chunk_text=chunk_text,
                            embedding=json.dumps(vec),
                        ))
                    except Exception as embed_err:
                        logger.warning(f"[UPLOAD] Embedding chunk {i} failed: {embed_err}")
                logger.info(f"[UPLOAD] Indexed {len(chunks)} chunks for article {art.id}")
            except Exception as re_idx_err:
                logger.warning(f"[UPLOAD] Indexing failed (best-effort): {re_idx_err}")

            db2.commit()
            logger.info(f"[UPLOAD] Enhanced: title={art.title!r} tags={bg_tags} errors={errs}")
        except Exception as e:
            logger.info(f"[UPLOAD] BG enhance failed: {e}")
            try:
                art = db2.query(Article).filter(Article.id == article_id).first()
                if art:
                    art.title = file.filename
                    art.content = raw_text + f"\n\n> ⚠️ 内容识别失败：{e}"
                    art.processing = None  # mark as done (failed, but no longer processing)
                    db2.commit()
            except Exception:
                pass
        finally:
            db2.close()

    asyncio.create_task(_bg_enhance())

    # Embeddings are computed after background recognition completes (see _bg_enhance)

    return article
